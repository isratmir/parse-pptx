import logging
import sys
import base64
import io

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx2geo import SvgMaker, ShapeToSVG
from PIL import Image

logging.basicConfig(stream=sys.stderr, level=logging.DEBUG)


def emuToPixels(emu):
    return round(emu / 9525) if emu is not None else emu


# prs = Presentation('presentation.pptx')
# prs = Presentation('mask.pptx')
prs = Presentation('face.pptx')

for num, slide in enumerate(prs.slides, start=1):
    # logging.debug(slide.slide_id)
    print(num)
    print("Slide has following objects:")
    print("\tSlide Id: ", slide.slide_id)
    print("\tSlide Open XML elements: ", slide.element)
    print("\tSlide layout: ", slide.slide_layout.name)
    print("\tSlide layout: ", slide.slide_layout)
    print("\tSlide background: ", slide.background.fill.type)
    if slide.background.fill.type == MSO_FILL.SOLID:
        print("\tSlide background RGB: ", slide.background.fill.fore_color)
    print("\tSlide placeholders: ")
    for placeholder in slide.placeholders:
        print("\t\tName: ", placeholder.shape_type)
        print("\t\tName: ", placeholder.name)
    print("\tShapes in slide:")
    for shape in slide.shapes:
        print("\t\tShape type: ", shape.shape_type)
        print("\t\t\tShape name: ", shape.name)
        print("\t\t\tShape rotation: ", shape.rotation)
        print("\t\t\tLeft (px): ", emuToPixels(shape.left))
        print("\t\t\tTop (px): ", emuToPixels(shape.top))
        print("\t\t\tWidth (px): ", emuToPixels(shape.width))
        print("\t\t\tHeight (px): ", emuToPixels(shape.height))
        print("\t\t\tWidth (emu): ", shape.width)
        print("\t\t\tHeight (emu): ", shape.height)
        if shape.is_placeholder:
            print("\t\t\tPlaceholder: ", shape)
            print("\t\t\tPlaceholder format: ", shape.placeholder_format.type)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print("\t\t\tPlaceholder text: ", run.text)
                    print("\t\t\t\tSize: ", emuToPixels(run.font.size))
                    print("\t\t\t\tFont name: ", run.font.name)
                    print("\t\t\t\tBold: ", run.font.bold)
                    print("\t\t\t\tItalic: ", run.font.italic)
                    print("\t\t\t\tColor type: ", run.font.color.type)
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print("\t\t\tRun: ", run.text)
                    print("\t\t\t\tSize: ", emuToPixels(run.font.size))
                    print("\t\t\t\tFont name: ", run.font.name)
                    print("\t\t\t\tBold: ", run.font.bold)
                    print("\t\t\t\tItalic: ", run.font.italic)
                    if run.font.color.type == MSO_COLOR_TYPE.RGB:
                        print("\t\t\t\tColor: ", run.font.color.rgb)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print("\t\t\tPicture name: ", shape.name)
            print("\t\t\tCrop top: ", emuToPixels(shape.crop_top))
            print("\t\t\tCrop bottom: ", emuToPixels(shape.crop_bottom))
            print("\t\t\tCrop left: ", emuToPixels(shape.crop_left))
            print("\t\t\tCrop right: ", emuToPixels(shape.crop_right))
            print("\t\t\tRotation: ", emuToPixels(shape.rotation))
            print("\t\t\tauto_shape_type: ", shape.auto_shape_type)
            print("\t\t\tFilename: ", shape.image.filename)
            print("\t\t\tFile size: ", shape.image.size)
            # print("\t\t\tImage ext: ", shape.image.ext)
            # print("\t\t\tImage ext: ", shape.image.content_type)
            # print("\t\t\tImage: ", shape.image.blob.decode('UTF-8'))

            i = Image.open(io.BytesIO(shape.image.blob))
            nuimage = i.resize((emuToPixels(shape.width), emuToPixels(shape.height)))
            nuimage.save('pictures/'+shape.name+'.'+shape.image.ext)

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            print("\t\t\tAuto name: ", shape.name)
            print("\t\t\tauto_shape_type: ", shape.auto_shape_type)
            if shape.line.color.type == MSO_COLOR_TYPE.RGB:
                print("\t\t\tline color: ", shape.line.color.rgb)
            if shape.line.color.type == MSO_COLOR_TYPE.SCHEME:
                print("\t\t\tline color: ", shape.line.color.theme_color)
            print("\t\t\tline width: ", emuToPixels(shape.line.width))
            print("\t\t\tfill type: ", shape.fill.type)
            if shape.fill.type == MSO_FILL.SOLID:
                if shape.fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                    print("\t\t\tfill color: ", shape.fill.fore_color.rgb)
                if shape.fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                    print("\t\t\tfill color: ", shape.fill.fore_color.theme_color)
            print("\t\t\tback color: ", shape.fill)
            print("\t\t\tshape_type: ", shape.shape_type)
            print("\t\t\tadjustments: ", shape.adjustments)
            for adjustment in shape.adjustments:
                print("Adj", adjustment)
            ShapeToSVG(shape)

    print("-----------------------------------------------------------")
    # exit()
